#!/usr/bin/env python3
"""Update slides 10-12 in 제로스트 2차 사전보고서.pptx."""

from __future__ import annotations

from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

PPT_PATH = Path("/Users/yoonnarae/Downloads/제로스트 2차 사전보고서.pptx")
GALLERY_PATH = Path("/Users/yoonnarae/Downloads/zerost-ppt-assets/slide11-fe-phone-gallery.png")

# Deck palette (matches existing slides)
GREEN = RGBColor(0x10, 0xB9, 0x81)
GREEN_LIGHT = RGBColor(0x34, 0xD3, 0x99)
GREEN_DARK = RGBColor(0x06, 0x4E, 0x3B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xD1, 0xD5, 0xDB)
GRAY_MUTED = RGBColor(0x9C, 0xA3, 0xAF)
BLACK = RGBColor(0x00, 0x00, 0x00)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
ORANGE = RGBColor(0xF9, 0x73, 0x16)


def emu_in(v: float):
    return Inches(v)


def clear_slide(slide) -> None:
    sp_tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        if shape.name.startswith("Picture") or shape.name.startswith("TextBox") or shape.name.startswith("Rounded"):
            sp = shape._element
            sp_tree.remove(sp)


def add_rect(slide, left, top, width, height, fill, radius=None):
    if radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    if radius:
        shape.adjustments[0] = 0.08
    return shape


def add_text(slide, left, top, width, height, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Pretendard"
    return box


def add_header(slide, badge: str, page: str, title: str, subtitle: str = "") -> None:
    add_rect(slide, emu_in(0), emu_in(0), emu_in(13.333), emu_in(7.5), GREEN_DARK)
    add_rect(slide, emu_in(0.8), emu_in(0.55), emu_in(3.2), emu_in(0.42), GREEN, radius=True)
    add_text(slide, emu_in(0.95), emu_in(0.58), emu_in(3.0), emu_in(0.35), badge, size=11, bold=True, color=WHITE)
    add_text(slide, emu_in(11.5), emu_in(0.58), emu_in(1.5), emu_in(0.35), page, size=14, bold=True, color=GREEN_LIGHT, align=PP_ALIGN.RIGHT)
    add_text(slide, emu_in(0.8), emu_in(1.15), emu_in(11.5), emu_in(0.8), title, size=48, bold=False, color=WHITE)
    if subtitle:
        add_text(slide, emu_in(0.8), emu_in(1.95), emu_in(11.5), emu_in(0.5), subtitle, size=20, color=GRAY)


def add_footer(slide) -> None:
    add_text(slide, emu_in(0.8), emu_in(6.95), emu_in(6), emu_in(0.35), "ZERO-WASTE HABIT LOOP : ZEROST", size=14, bold=True, color=GREEN)


def update_page_numbers(prs) -> None:
    mapping = {
        0: "1 / 12",
        1: "2 / 12",
        2: "3 / 12",
        3: "4 / 12",
        4: "5 / 12",
        5: "6 / 12",
        6: "7 / 12",
        7: "8 / 12",
        8: "9 / 12",
    }
    for idx, new_page in mapping.items():
        slide = prs.slides[idx]
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip().endswith("/ 10"):
                shape.text_frame.text = new_page


def build_slide_10(slide) -> None:
    clear_slide(slide)
    add_header(
        slide,
        "FRONTEND DEMO",
        "10 / 12",
        "Frontend 화면 설계 (Demo)",
        "앱인토스 MVP · Granite · TDS · 18 screens",
    )
    add_rect(slide, emu_in(0.8), emu_in(2.7), emu_in(11.7), emu_in(3.8), RGBColor(0x02, 0x2C, 0x22), radius=True)
    bullets = [
        "온보딩 → 샵 선택 → 홈 · 미션 · 출석 · 스프 제작 · 가챠",
        "사진 1장 인증 · 에코잼 보상 · 알맹 포인트 연동 (파일럿)",
        "다음 슬라이드: 전체 18화면 스마트폰 갤러리",
    ]
    y = 3.0
    for b in bullets:
        add_text(slide, emu_in(1.2), emu_in(y), emu_in(10.8), emu_in(0.45), f"• {b}", size=18, color=GRAY)
        y += 0.55
    add_footer(slide)


def build_slide_11(slide) -> None:
    clear_slide(slide)
    add_header(
        slide,
        "FRONTEND DEMO",
        "11 / 12",
        "전체 화면 갤러리",
        "제로스트 FE mock — 모든 페이지 한눈에",
    )
    # Fit gallery below header
    left = emu_in(0.35)
    top = emu_in(2.45)
    width = emu_in(12.65)
    height = emu_in(4.35)
    slide.shapes.add_picture(str(GALLERY_PATH), left, top, width=width, height=height)
    add_footer(slide)


def add_check_card(slide, left, top, width, height, title, body, accent=GREEN):
    add_rect(slide, left, top, width, height, RGBColor(0x02, 0x2C, 0x22), radius=True)
    add_rect(slide, left + emu_in(0.15), top + emu_in(0.15), emu_in(0.55), emu_in(0.55), accent, radius=True)
    add_text(slide, left + emu_in(0.22), top + emu_in(0.2), emu_in(0.4), emu_in(0.4), "✓", size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, left + emu_in(0.85), top + emu_in(0.18), width - emu_in(1.0), emu_in(0.35), title, size=16, bold=True, color=GREEN_LIGHT)
    add_text(slide, left + emu_in(0.85), top + emu_in(0.55), width - emu_in(1.0), height - emu_in(0.65), body, size=13, color=GRAY)


def build_slide_12(slide) -> None:
    clear_slide(slide)
    add_header(
        slide,
        "PARTNERSHIP",
        "12 / 12",
        "알맹상점 파일럿 협의 결과",
        "2026.06.27 · 고금숙 펠로우 회신 · 협력 확정",
    )

    cards = [
        (
            "파일럿 매장 참여",
            "알맹상점을 파일럿 매장으로 확정.\n단골 고객 5~10명 모집·추천 진행 (SNS·매장 안내 등)",
            GREEN,
        ),
        (
            "참여 독려 & 소통",
            "테스트 기간 중 참여 독려 1~2회 가능.\n테스터 단톡방 개설·기간 중 소통 협의 가능",
            BLUE,
        ),
        (
            "포인트 · 리워드 연계",
            "알맹 포인트 매장 지급 가능 (매장 부담).\n토스 프런트 사용 중 → 앱인토스 연동 기반 마련",
            ORANGE,
        ),
        (
            "확장 · 추가 논의",
            "다른 제로웨이스트 샵 참여 문의 가능.\n환경실천 강화·포인트 지급 방식 추가 협의 환영",
            GREEN,
        ),
    ]

    x0 = 0.8
    y0 = 2.55
    w = 5.65
    h = 1.55
    gap_x = 0.35
    gap_y = 0.25
    for i, (title, body, accent) in enumerate(cards):
        col = i % 2
        row = i // 2
        left = emu_in(x0 + col * (w + gap_x))
        top = emu_in(y0 + row * (h + gap_y))
        add_check_card(slide, left, top, emu_in(w), emu_in(h), title, body, accent)

    add_rect(slide, emu_in(0.8), emu_in(6.05), emu_in(11.7), emu_in(0.75), RGBColor(0x02, 0x2C, 0x22), radius=True)
    add_text(
        slide,
        emu_in(1.05),
        emu_in(6.18),
        emu_in(11.2),
        emu_in(0.5),
        "파일럿 일정 (조율): 7월 중순~말 · 1주 내외 · 사전·사후 구글폼 · 감사 상품 (예산 30만 원)",
        size=14,
        color=GRAY,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide)


def main() -> None:
    if not PPT_PATH.exists():
        raise SystemExit(f"PPT not found: {PPT_PATH}")
    if not GALLERY_PATH.exists():
        raise SystemExit(f"Gallery not found: {GALLERY_PATH}")

    prs = Presentation(str(PPT_PATH))
    update_page_numbers(prs)
    build_slide_10(prs.slides[9])
    build_slide_11(prs.slides[10])
    build_slide_12(prs.slides[11])
    prs.save(str(PPT_PATH))
    print(f"Updated: {PPT_PATH}")


if __name__ == "__main__":
    main()
